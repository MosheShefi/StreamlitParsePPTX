# Import convention
import streamlit as st
from pptx import Presentation
from pretty_notification_box import notification_box
import pandas as pd
# from streamlit_custom_notification_box import custom_notification_box
# import glob
# import pathlib
# import os.path

# st.subheader("Component with constant args")

styles = {'material-icons':{'color': 'red'},
          'title': {'font-weight':'bold'},
          'notification-content-container': {'':''},
          'title-text-url-container': {'',''},
          'notification-text-link-close-container': {'',''},
          'external-link': {'',''},
          'close-button': {'',''}}

uploaded_files = st.file_uploader("Choose a PPTX file",type=["pptx"],
                                  accept_multiple_files=True)

def upload():
    if not uploaded_files:
        st.text('Load pptx file/s first')
    else:
        for uploaded_file in uploaded_files:
            prs = Presentation(uploaded_file)
            st.write("---------- " + uploaded_file.name + " -------------")
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        st.write(shape.text)

st.button("Parse text out of pptx file/s", on_click=upload, disabled=not uploaded_files)

